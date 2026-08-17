"""Generic index slide reflow and slide-number updates."""

from __future__ import annotations

import copy
import logging
import re

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from rapidfuzz import fuzz

from app.services.ppt_shape_utils import normalize_title_text, slide_title_text
from app.wsr_engine.models import ProjectMap

logger = logging.getLogger(__name__)

_INDEX_NUMBER_RE = re.compile(r"^\d{1,2}\u200b?$")

_INDEX_LABEL_KEYWORDS: tuple[str, ...] = (
    "cost core",
    "supplier core",
    "pricing core",
    "wentworth",
    "location core",
    "pharmacy",
    "wellness",
    "global sourcing",
    "product attribute",
    "loco",
    "bsa",
    "matters of attention",
    "team allocation",
)


def _normalize_index_text(text: str) -> str:
    return normalize_title_text(text).lower()


def _find_index_slide_index(prs: Presentation) -> int | None:
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _normalize_index_text(shape.text_frame.text) == "index":
                return i
    return None


def _find_index_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def _index_table_cells_row_major(table):
    cells = []
    for row in table.rows:
        cells.extend(row.cells)
    return cells


def _cell_has_index_content(cell) -> bool:
    text = _normalize_index_text(cell.text_frame.text)
    if not text:
        return False
    if any(kw in text for kw in _INDEX_LABEL_KEYWORDS):
        return True
    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            return True
    return False


def _is_index_number_paragraph(paragraph) -> bool:
    text = normalize_title_text("".join(run.text for run in paragraph.runs))
    return bool(_INDEX_NUMBER_RE.match(text))


def _set_paragraph_slide_number(paragraph, slide_number: int) -> None:
    display = f"{slide_number:02d}"
    if paragraph.runs:
        paragraph.runs[0].text = display
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = display


def _set_element_hyperlink_to_slide(hlink_elem, index_slide, target_slide_idx, prs) -> None:
    target_part = prs.slides[target_slide_idx].part
    slide_rid = index_slide.part.relate_to(target_part, RT.SLIDE)
    hlink_elem.set(qn("r:id"), slide_rid)


def _update_cell_index_entry(cell, index_slide, target_slide_idx, prs) -> bool:
    changed = False
    display_num = target_slide_idx + 1

    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            old = normalize_title_text("".join(run.text for run in paragraph.runs))
            new = f"{display_num:02d}"
            if old != new:
                _set_paragraph_slide_number(paragraph, display_num)
                changed = True

    tx_body = cell.text_frame._txBody
    for p_elem in tx_body.findall(qn("a:p")):
        for hlink in p_elem.iter(qn("a:hlinkClick")):
            _set_element_hyperlink_to_slide(hlink, index_slide, target_slide_idx, prs)
            changed = True

    return changed


def _clear_index_cell_completely(cell) -> None:
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""
    tx_body = cell.text_frame._txBody
    for p_elem in tx_body.findall(qn("a:p")):
        for hlink in list(p_elem.findall(qn("a:hlinkClick"))):
            p_elem.remove(hlink)


def _clone_cell_text_body(dst_cell, src_tx_body) -> None:
    dst_tc = dst_cell._tc
    old_tx = dst_tc.find(qn("a:txBody"))
    if old_tx is not None:
        dst_tc.remove(old_tx)
    dst_tc.insert(0, copy.deepcopy(src_tx_body))


def _resolve_index_target(
    prs: Presentation,
    cell_text: str,
    projects: list[ProjectMap],
) -> int | None:
    label = _normalize_index_text(cell_text)
    if not label:
        return None

    # Skip non-delivery index rows (team allocation, RAID, etc.)
    if "team allocation" in label or "raid status" in label:
        return None

    best_score = 0
    best_idx: int | None = None

    for proj in projects:
        score = fuzz.partial_ratio(label, proj.project_name.lower())
        if score > best_score and score >= 70:
            best_score = score
            best_idx = proj.main_slide_index

    if best_idx is not None and best_idx < len(prs.slides):
        return best_idx

    for i, slide in enumerate(prs.slides):
        title = normalize_title_text(slide_title_text(slide)).lower()
        if not title or "(contd" in title:
            continue
        score = fuzz.partial_ratio(label, title)
        if score > best_score and score >= 70:
            best_score = score
            best_idx = i

    return best_idx


def _collect_active_index_entries(prs: Presentation, table, projects: list[ProjectMap]) -> list[tuple[int, object]]:
    entries: list[tuple[int, object]] = []
    seen_targets: set[int] = set()

    for cell in _index_table_cells_row_major(table):
        if not _cell_has_index_content(cell):
            continue
        cell_text = _normalize_index_text(cell.text_frame.text)
        target_idx = _resolve_index_target(prs, cell_text, projects)
        if target_idx is None or target_idx in seen_targets:
            continue
        seen_targets.add(target_idx)
        entries.append((target_idx, copy.deepcopy(cell.text_frame._txBody)))

    return entries


def reflow_index_slide(prs: Presentation, projects: list[ProjectMap]) -> int:
    index_idx = _find_index_slide_index(prs)
    if index_idx is None:
        logger.warning("Index slide not found")
        return 0

    index_slide = prs.slides[index_idx]
    table = _find_index_table(index_slide)
    if table is None:
        logger.warning("Index table not found")
        return 0

    active_entries = _collect_active_index_entries(prs, table, projects)
    slot_cells = _index_table_cells_row_major(table)

    for cell in slot_cells:
        if _cell_has_index_content(cell):
            _clear_index_cell_completely(cell)

    updated = 0
    for slot_idx, (target_idx, tx_body) in enumerate(active_entries):
        if slot_idx >= len(slot_cells):
            break
        cell = slot_cells[slot_idx]
        _clone_cell_text_body(cell, tx_body)
        if _update_cell_index_entry(cell, index_slide, target_idx, prs):
            updated += 1

    logger.info("Index reflow: %d entries updated", updated)
    return updated
