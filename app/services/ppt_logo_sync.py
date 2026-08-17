"""Re-insert header logos on continuation slides after shape XML copy."""

from __future__ import annotations

import io
import logging

from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

HEB_LOGO_SHAPE_ID = 3


def _delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _picture_has_blob(shape) -> bool:
    try:
        _ = shape.image.blob
        return True
    except (AttributeError, ValueError, KeyError):
        return False


def find_heb_logo_on_slide(slide):
    """Return the H-E-B header logo picture on a main project slide."""
    candidate = next(
        (s for s in slide.shapes if getattr(s, "shape_id", None) == HEB_LOGO_SHAPE_ID),
        None,
    )
    if candidate is None:
        return None
    if candidate.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    if not _picture_has_blob(candidate):
        return None
    return candidate


def sync_heb_logo_from_main(main_slide, contd_slide) -> bool:
    """
    Copy the H-E-B header logo onto a (Contd…) slide.

    ``copy_shapes_to_slide`` deep-copies shape XML but not image part
    relationships, which leaves a broken picture placeholder on contd slides.
    Re-insert the logo from the main slide's embedded image bytes.
    """
    src = find_heb_logo_on_slide(main_slide)
    if src is None:
        logger.warning("HEB logo (shape id %s) not found on main slide", HEB_LOGO_SHAPE_ID)
        return False

    blob = src.image.blob
    for sh in list(contd_slide.shapes):
        sid = getattr(sh, "shape_id", None)
        if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if sid == HEB_LOGO_SHAPE_ID or not _picture_has_blob(sh):
            _delete_shape(sh)

    contd_slide.shapes.add_picture(
        io.BytesIO(blob), src.left, src.top, width=src.width, height=src.height
    )
    return True
